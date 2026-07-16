# -*- coding: utf-8 -*-
"""
deck_helpers.py — a small python-pptx "design system" for building clean,
consistent 16:9 slide decks as code.

Usage:
    from deck_helpers import *
    prs = new_deck()
    s = slide(prs)
    title(s, "页面标题", "副标题", chip="Tag", hero="hero_img")   # hero optional
    txt(s, 0.6, 2.0, 6, 1, "正文，**加粗片段**用双星号，\\n 换行", size=14)
    card(s, 0.6, 3.0, 4, 1.5, fill=G_BLUE50, line=None)
    pic(s, "anker_images/p08_nas.png", 7, 2, 5, 3)              # crop-to-fill
    save(prs, "out.pptx")

Conventions
- Canvas is 13.333" x 7.5" (16:9). All positions are absolute, in inches.
- `txt` is the workhorse: `**...**` toggles bold; `\\n` makes new paragraphs.
- `pic` crops to fill the target box (no distortion); falls back to a grey
  placeholder if the file is missing.
- Colors are the Google palette; tweak to match any brand.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ---- palette (swap for your brand) ----
G_BLUE  = RGBColor(0x42, 0x85, 0xF4); G_RED    = RGBColor(0xEA, 0x43, 0x35)
G_YELLOW= RGBColor(0xFB, 0xBC, 0x04); G_GREEN  = RGBColor(0x34, 0xA8, 0x53)
G_GREY9 = RGBColor(0x20, 0x21, 0x24); G_GREY7  = RGBColor(0x5F, 0x63, 0x68)
G_GREY3 = RGBColor(0xDA, 0xDC, 0xE0); G_GREY1  = RGBColor(0xF8, 0xF9, 0xFA)
G_BLUE50= RGBColor(0xE8, 0xF0, 0xFE); G_YEL50  = RGBColor(0xFE, 0xF7, 0xE0)
G_GRN50 = RGBColor(0xE6, 0xF4, 0xEA); G_RED50  = RGBColor(0xFC, 0xE8, 0xE6)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF); DARK     = RGBColor(0x1A, 0x1C, 0x1E)
FONT = "Google Sans"          # any installed font; falls back gracefully
IMG  = "."                    # base dir for images used by pic()/title(hero=)

def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    return prs

def slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # blank

def save(prs, path):
    prs.save(path)
    print(f"saved {path} · {len(prs.slides._sldIdLst)} slides")

# ---------- text ----------
def txt(s, x, y, w, h, text, size=14, bold=False, color=G_GREY9,
        align=PP_ALIGN.LEFT, spacing=1.0, anchor=MSO_ANCHOR.TOP):
    """Textbox. '**' toggles bold inside a line; '\\n' starts a new paragraph."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        for j, seg in enumerate(line.split("**")):
            r = p.add_run(); r.text = seg
            r.font.size = Pt(size); r.font.name = FONT
            r.font.color.rgb = color; r.font.bold = bold or (j % 2 == 1)
    return tb

# ---------- shapes ----------
def card(s, x, y, w, h, fill=WHITE, line=G_GREY3, r=0.08):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    try: sp.adjustments[0] = r
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line: sp.line.color.rgb = line; sp.line.width = Pt(1)
    else: sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def bar(s, x, y, w, h, c):
    """Solid rectangle. Use as accent strips, thin lines (tiny w/h), fills."""
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = c
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def hline(s, x, y, w, c=G_GREY3, pt=1):
    l = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                               Inches(x), Inches(y), Inches(x + w), Inches(y))
    l.line.color.rgb = c; l.line.width = Pt(pt)
    return l

def dots(s, x=0.6, y=0.45):
    """The Google four-dot motif (top-left brand mark)."""
    for i, c in enumerate([G_BLUE, G_RED, G_YELLOW, G_GREEN]):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + i*0.22), Inches(y),
                               Inches(0.12), Inches(0.12))
        d.fill.solid(); d.fill.fore_color.rgb = c
        d.line.fill.background(); d.shadow.inherit = False

# ---------- media ----------
def pic(s, name, x, y, w, h, back=False):
    """Image cropped to FILL the box (keeps aspect, no distortion).
    `name` is relative to IMG unless it already ends in an image extension path.
    Missing file -> grey placeholder card."""
    p = f"{IMG}/{name}" if not os.path.isabs(name) else name
    if not os.path.exists(p) and not name.lower().endswith((".png", ".jpg", ".jpeg")):
        p = f"{IMG}/{name}.png"
    if not os.path.exists(p):
        card(s, x, y, w, h, fill=G_GREY1); return None
    from PIL import Image
    iw, ih = Image.open(p).size
    tr, ir = w / h, iw / ih
    img = s.shapes.add_picture(p, Inches(x), Inches(y), Inches(w), Inches(h))
    if ir > tr:                       # image wider -> crop sides
        c = (1 - tr/ir) / 2; img.crop_left = c; img.crop_right = c
    else:                             # image taller -> crop top/bottom
        c = (1 - ir/tr) / 2; img.crop_top = c; img.crop_bottom = c
    if back:                          # send picture behind other shapes
        t = s.shapes._spTree; t.remove(img._element); t.insert(2, img._element)
    return img

def video(s, name, x, y, w, h, poster=None, mime="video/mp4"):
    p = f"{IMG}/{name}"; pp = f"{IMG}/{poster}" if poster else None
    if not os.path.exists(p):
        card(s, x, y, w, h, fill=G_GREY1)
        txt(s, x, y+h/2-0.2, w, 0.4, f"▶ {name}", size=11,
            color=G_GREY7, align=PP_ALIGN.CENTER); return
    s.shapes.add_movie(p, Inches(x), Inches(y), Inches(w), Inches(h),
                       poster_frame_image=pp, mime_type=mime)

def scrim(s, x, y, w, h, alpha=55):
    """Semi-transparent dark overlay (for text legibility over a hero image)."""
    sp = bar(s, x, y, w, h, DARK)
    from pptx.oxml.ns import qn
    clr = sp.fill._xPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    clr.append(clr.makeelement(qn('a:alpha'), {'val': str(alpha*1000)}))
    return sp

# ---------- composite ----------
def title(s, main, sub=None, chip=None, chip_color=G_BLUE, hero=None):
    """Page title. With hero=<image>, draws a 21:9 hero strip + dark scrim and
    white title text; otherwise a clean white title with an accent underline."""
    if hero:
        pic(s, hero, 0, 0, 13.333, 1.7); scrim(s, 0, 0, 13.333, 1.7, 55)
        dots(s, 0.6, 0.35)
        sz = 26 if len(main) <= 24 else (22 if len(main) <= 34 else 19)
        txt(s, 0.6, 0.62, 11.5, 0.6, main, size=sz, bold=True, color=WHITE)
        if sub: txt(s, 0.6, 1.2, 11.5, 0.35, sub, size=13,
                    color=RGBColor(0xE0, 0xE0, 0xE0))
        if chip:
            cw = max(1.8, len(chip)*0.15 + 0.5)
            card(s, 12.7-cw, 0.55, cw, 0.45, fill=WHITE, line=None)
            txt(s, 12.7-cw, 0.62, cw, 0.32, chip, size=11, bold=True,
                color=chip_color, align=PP_ALIGN.CENTER)
    else:
        dots(s); sz = 28 if len(main) <= 22 else (24 if len(main) <= 32 else 21)
        txt(s, 0.6, 0.72, 11.5, 0.7, main, size=sz, bold=True)
        bar(s, 0.6, 1.42, 0.6, 0.06, G_BLUE)
        if sub: txt(s, 1.45, 1.32, 10.5, 0.35, sub, size=13, color=G_GREY7)
        if chip:
            cw = max(1.8, len(chip)*0.15 + 0.5)
            card(s, 12.7-cw, 0.72, cw, 0.42, fill=G_BLUE50, line=None)
            txt(s, 12.7-cw, 0.77, cw, 0.32, chip, size=11, bold=True,
                color=chip_color, align=PP_ALIGN.CENTER)

def notes(s, t):
    s.notes_slide.notes_text_frame.text = t

def kpi(s, x, y, v, label, sub, c, w=2.85, h=1.55):
    """A KPI card: big number + label + sub. Lay several in a row."""
    card(s, x, y, w, h)
    txt(s, x+0.2, y+0.15, w-0.4, 0.5, v, size=22, bold=True, color=c)
    txt(s, x+0.2, y+0.72, w-0.4, 0.3, label, size=12, bold=True)
    txt(s, x+0.2, y+1.05, w-0.4, 0.4, sub, size=10, color=G_GREY7)

def callout(s, x, y, w, text, fill=G_BLUE50, bar_c=G_BLUE, h=0.55, size=12):
    """Bottom 'so what' band: tinted box + left accent bar + bold-aware text."""
    card(s, x, y, w, h, fill=fill, line=None)
    bar(s, x, y, 0.08, h, bar_c)
    txt(s, x+0.3, y+0.08, w-0.5, h-0.1, text, size=size, spacing=1.2)
